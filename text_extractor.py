import os
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import logging

class TextExtractor:
    """
    A comprehensive text extractor that can handle multiple file types:
    - PDF files
    - Word documents (.docx)
    - PowerPoint presentations (.pptx)
    - Excel spreadsheets (.xlsx)
    """
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_from_pdf,
            '.docx': self._extract_from_docx,
            '.pptx': self._extract_from_pptx,
            '.xlsx': self._extract_from_xlsx,
            '.xls': self._extract_from_xlsx
        }
    
    def extract_text(self, file_path):
        """
        Extract text from a file based on its extension.
        
        Args:
            file_path (str): Path to the file
            
        Returns:
            tuple: (extracted_text, metadata_dict) or (None, error_message)
        """
        try:
            if not os.path.exists(file_path):
                return None, "File not found"
            
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension not in self.supported_formats:
                return None, f"Unsupported file format: {file_extension}"
            
            # Extract text using the appropriate method
            extractor_method = self.supported_formats[file_extension]
            return extractor_method(file_path)
            
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {str(e)}")
            return None, f"Error processing file: {str(e)}"
    
    def _extract_from_pdf(self, file_path):
        """
        Extract text from PDF files.
        """
        try:
            with open(file_path, 'rb') as pdfFileObj:
                pdfReader = PyPDF2.PdfReader(pdfFileObj)
                num_pages = len(pdfReader.pages)
                
                text = ""
                for i in range(num_pages):
                    pageObj = pdfReader.pages[i]
                    page_text = pageObj.extract_text()
                    text += page_text + "\n"
                
                metadata = {
                    'file_type': 'PDF',
                    'pages': num_pages,
                    'characters': len(text),
                    'words': len(text.split())
                }
                
                return text.strip(), metadata
                
        except Exception as e:
            return None, f"PDF extraction error: {str(e)}"
    
    def _extract_from_docx(self, file_path):
        """
        Extract text from Word documents (.docx).
        """
        try:
            doc = Document(file_path)
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\t"
                    text += "\n"
            
            metadata = {
                'file_type': 'Word Document',
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'characters': len(text),
                'words': len(text.split())
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            return None, f"Word document extraction error: {str(e)}"
    
    def _extract_from_pptx(self, file_path):
        """
        Extract text from PowerPoint presentations (.pptx).
        """
        try:
            prs = Presentation(file_path)
            text = ""
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                text += f"\n--- Slide {slide_count} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text += cell.text + "\t"
                            text += "\n"
            
            metadata = {
                'file_type': 'PowerPoint Presentation',
                'slides': slide_count,
                'characters': len(text),
                'words': len(text.split())
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            return None, f"PowerPoint extraction error: {str(e)}"
    
    def _extract_from_xlsx(self, file_path):
        """
        Extract text from Excel spreadsheets (.xlsx, .xls).
        """
        try:
            wb = load_workbook(file_path, data_only=True)
            text = ""
            sheet_count = 0
            
            for sheet_name in wb.sheetnames:
                sheet_count += 1
                ws = wb[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Get the maximum row and column numbers
                max_row = ws.max_row
                max_col = ws.max_column
                
                for row in range(1, max_row + 1):
                    row_text = ""
                    for col in range(1, max_col + 1):
                        cell_value = ws.cell(row=row, column=col).value
                        if cell_value is not None:
                            row_text += str(cell_value) + "\t"
                    if row_text.strip():
                        text += row_text.strip() + "\n"
            
            metadata = {
                'file_type': 'Excel Spreadsheet',
                'sheets': sheet_count,
                'characters': len(text),
                'words': len(text.split())
            }
            
            return text.strip(), metadata
            
        except Exception as e:
            return None, f"Excel extraction error: {str(e)}"
    
    def get_supported_formats(self):
        """
        Get list of supported file formats.
        """
        return list(self.supported_formats.keys())
    
    def is_supported(self, file_path):
        """
        Check if a file format is supported.
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        return file_extension in self.supported_formats 